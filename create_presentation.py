"""
Clean presentation with illustrations for the Data Chat Agent project.
Minimal typography + generated visuals on key slides.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Palette ──────────────────────────────────────────────────────────────
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE  = RGBColor(0xF8, 0xF9, 0xFB)
NEAR_BLACK = RGBColor(0x1A, 0x1A, 0x2E)
DARK       = RGBColor(0x12, 0x14, 0x26)
GRAY       = RGBColor(0x71, 0x71, 0x80)
LIGHT_GRAY = RGBColor(0xA0, 0xA0, 0xAE)
RULE_GRAY  = RGBColor(0xE0, 0xE0, 0xE5)
BLUE       = RGBColor(0x33, 0x6B, 0xF5)
SOFT_BLUE  = RGBColor(0xEB, 0xF0, 0xFE)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
ASSETS  = os.path.join(os.path.dirname(os.path.abspath(__file__)), "presentation_assets")

# ── Helpers ──────────────────────────────────────────────────────────────
def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def txt(slide, left, top, w, h, text, size=18, color=NEAR_BLACK,
        bold=False, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.font.italic = italic
    p.alignment = align
    return box

def thin_line(slide, left, top, width, color=RULE_GRAY):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def accent_dot(slide, left, top, color=BLUE):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Inches(0.12), Inches(0.12))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()

def add_img(slide, path, left, top, width=None, height=None):
    """Add an image; at least width or height must be given."""
    if width:
        slide.shapes.add_picture(path, left, top, width=width)
    else:
        slide.shapes.add_picture(path, left, top, height=height)

def numbered_item(slide, left, top, num, title, desc, w=Inches(5)):
    txt(slide, left, top, Inches(0.5), Inches(0.45),
        num, size=28, color=BLUE, bold=True)
    txt(slide, left + Inches(0.55), top + Inches(0.02), w, Inches(0.4),
        title, size=17, color=NEAR_BLACK, bold=True)
    txt(slide, left + Inches(0.55), top + Inches(0.38), w, Inches(0.5),
        desc, size=13, color=GRAY)

# ── Build ────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]

LM = Inches(1.4)
TM = Inches(0.9)

# =====================================================================
# 1  TITLE  — big headline + hero illustration on the right
# =====================================================================
s = prs.slides.add_slide(blank)
bg(s, WHITE)

txt(s, LM, Inches(1.8), Inches(6), Inches(1.2),
    "Data Chat Agent", size=52, color=NEAR_BLACK, bold=True)

txt(s, LM, Inches(3.1), Inches(5.5), Inches(0.8),
    "Talk to any dataset in plain English.\nGet answers, charts, and insights — instantly.",
    size=22, color=GRAY)

thin_line(s, LM, Inches(4.3), Inches(3.5), BLUE)

txt(s, LM, Inches(4.65), Inches(6), Inches(0.4),
    "Python  ·  FastAPI  ·  Next.js  ·  Google Gemini  ·  Vega-Lite",
    size=13, color=LIGHT_GRAY)

txt(s, LM, Inches(6.3), Inches(5), Inches(0.35),
    "Project Demo  —  March 2026", size=13, color=LIGHT_GRAY)

# Hero illustration on right
add_img(s, os.path.join(ASSETS, "hero.png"),
        Inches(7.5), Inches(0.8), width=Inches(5.3))


# =====================================================================
# 2  WHAT IS IT?  — text left, illustration right
# =====================================================================
s = prs.slides.add_slide(blank)
bg(s, WHITE)

txt(s, LM, TM, Inches(4), Inches(0.35),
    "What is it?", size=13, color=BLUE, bold=True)

txt(s, LM, Inches(1.6), Inches(6.5), Inches(1.6),
    "An AI chatbot that lets you upload any CSV or Excel file and explore it using natural language.",
    size=28, color=NEAR_BLACK, bold=True)

thin_line(s, LM, Inches(3.5), Inches(6))

# Three short points
y = Inches(4.0)
points = [
    ("Upload your data",
     "Drop in any .csv or .xlsx — schema and column types are detected automatically."),
    ("Ask questions",
     "\"Show top 5 by price\"  ·  \"Plot average salary by dept\"  ·  \"How many per category?\""),
    ("Get answers + charts",
     "The agent writes SQL, runs it, and returns text summaries or interactive Vega-Lite charts."),
]
for title, desc in points:
    accent_dot(s, LM, y + Inches(0.06))
    txt(s, LM + Inches(0.3), y, Inches(5.8), Inches(0.35),
        title, size=15, color=NEAR_BLACK, bold=True)
    txt(s, LM + Inches(0.3), y + Inches(0.35), Inches(5.8), Inches(0.45),
        desc, size=12, color=GRAY)
    y += Inches(0.95)

# Illustration right side
add_img(s, os.path.join(ASSETS, "query.png"),
        Inches(8.4), Inches(1.5), width=Inches(4.2))


# =====================================================================
# 3  THE PROBLEM
# =====================================================================
s = prs.slides.add_slide(blank)
bg(s, WHITE)

txt(s, LM, TM, Inches(4), Inches(0.35),
    "The problem", size=13, color=BLUE, bold=True)

txt(s, LM, Inches(1.6), Inches(10), Inches(1.0),
    "Business users have data but not the skills to explore it.",
    size=30, color=NEAR_BLACK, bold=True)

y = Inches(3.2)
problems = [
    ("People depend on analysts for every ad-hoc question.",
     "A simple \"how many X by Y\" can take hours when you need someone else to write the query."),
    ("Dashboards are rigid — they can't answer new questions.",
     "Pre-built BI tools only cover known questions. Anything new requires a ticket and a wait."),
    ("Most tools are tied to one dataset.",
     "Switching to a new file means rebuilding schemas, queries, and visualizations from scratch."),
]
for title, desc in problems:
    accent_dot(s, LM, y + Inches(0.06))
    txt(s, LM + Inches(0.3), y, Inches(10), Inches(0.35),
        title, size=16, color=NEAR_BLACK, bold=True)
    txt(s, LM + Inches(0.3), y + Inches(0.4), Inches(9), Inches(0.4),
        desc, size=13, color=GRAY)
    y += Inches(1.15)


# =====================================================================
# 4  HOW IT WORKS — steps + illustration
# =====================================================================
s = prs.slides.add_slide(blank)
bg(s, WHITE)

txt(s, LM, TM, Inches(4), Inches(0.35),
    "How it works", size=13, color=BLUE, bold=True)

txt(s, LM, Inches(1.6), Inches(10), Inches(0.8),
    "Five steps from file to insight.", size=30, color=NEAR_BLACK, bold=True)

thin_line(s, LM, Inches(2.8), Inches(10.5))

steps = [
    ("01", "Upload",  "CSV or Excel file\nvia the sidebar."),
    ("02", "Detect",  "Auto-create SQLite\ntable and metadata."),
    ("03", "Ask",     "Type a question in\nplain English."),
    ("04", "Process", "Gemini picks the tool\nand generates SQL."),
    ("05", "Respond", "Text summary or\ninteractive chart."),
]
col_w = Inches(2.0)
y = Inches(3.2)
for i, (num, title, desc) in enumerate(steps):
    x = LM + i * (col_w + Inches(0.15))
    txt(s, x, y, col_w, Inches(0.55), num, size=36, color=BLUE, bold=True)
    txt(s, x, y + Inches(0.6), col_w, Inches(0.35),
        title, size=17, color=NEAR_BLACK, bold=True)
    txt(s, x, y + Inches(1.0), col_w, Inches(0.9), desc, size=12, color=GRAY)


# =====================================================================
# 5  ARCHITECTURE — illustration-driven
# =====================================================================
s = prs.slides.add_slide(blank)
bg(s, WHITE)

txt(s, LM, TM, Inches(4), Inches(0.35),
    "Architecture", size=13, color=BLUE, bold=True)

txt(s, LM, Inches(1.6), Inches(5.5), Inches(0.8),
    "Four components, loosely coupled.",
    size=30, color=NEAR_BLACK, bold=True)

# Architecture illustration — large, centered
add_img(s, os.path.join(ASSETS, "architecture.png"),
        Inches(7.2), Inches(0.9), width=Inches(5.5))

# Simple text list of components below on left
y = Inches(3.0)
components = [
    ("Frontend",  "Next.js 14, React, Tailwind CSS, Vega-Lite charts"),
    ("Backend",   "FastAPI, Python — chat agent with tool router & SQL builder"),
    ("LLM",       "Google Gemini 2.5 Flash — intent detection & text-to-SQL"),
    ("Database",  "SQLite — auto-created per uploaded dataset"),
]
for label, desc in components:
    accent_dot(s, LM, y + Inches(0.06))
    txt(s, LM + Inches(0.28), y, Inches(1.5), Inches(0.35),
        label, size=15, color=NEAR_BLACK, bold=True)
    txt(s, LM + Inches(1.85), y, Inches(4.5), Inches(0.35),
        desc, size=13, color=GRAY)
    y += Inches(0.7)

# Endpoint count callout
txt(s, LM, Inches(6.2), Inches(6), Inches(0.3),
    "8 REST API endpoints  ·  2 tool types (data_query, data_stats)  ·  4 chart types",
    size=12, color=LIGHT_GRAY)


# =====================================================================
# 6  TECH STACK — clean table
# =====================================================================
s = prs.slides.add_slide(blank)
bg(s, WHITE)

txt(s, LM, TM, Inches(4), Inches(0.35),
    "Tech stack", size=13, color=BLUE, bold=True)

left_x = LM
right_x = LM + Inches(3.2)
y = Inches(1.8)
row_h = Inches(0.95)

stack = [
    ("AI Model",    "Google Gemini 2.5 Flash via LangChain"),
    ("Backend",     "Python 3.10+, FastAPI, Uvicorn"),
    ("Database",    "SQLite — auto-created from uploaded files"),
    ("Data Layer",  "Pandas for ingestion, Pydantic for validation"),
    ("Frontend",    "Next.js 14, React 18, TypeScript, Tailwind CSS"),
    ("Charts",      "Vega-Lite specs built server-side, rendered via vega-embed"),
]
for label, detail in stack:
    thin_line(s, left_x, y, Inches(10))
    txt(s, left_x, y + Inches(0.18), Inches(3), Inches(0.35),
        label, size=15, color=NEAR_BLACK, bold=True)
    txt(s, right_x, y + Inches(0.18), Inches(7), Inches(0.35),
        detail, size=15, color=GRAY)
    y += row_h
thin_line(s, left_x, y, Inches(10))


# =====================================================================
# 7  KEY FEATURES
# =====================================================================
s = prs.slides.add_slide(blank)
bg(s, WHITE)

txt(s, LM, TM, Inches(4), Inches(0.35),
    "Key features", size=13, color=BLUE, bold=True)

txt(s, LM, Inches(1.6), Inches(10), Inches(0.8),
    "What makes it useful.", size=30, color=NEAR_BLACK, bold=True)

y = Inches(3.0)
features = [
    ("Dataset agnostic",       "Works with any CSV or Excel — no hardcoded schemas."),
    ("Text-to-SQL",            "Plain English converted to safe, parameterized SQL."),
    ("Smart charting",         "Picks the right chart type from your question."),
    ("Auto-generated context", "Column types, ranges, and samples detected on upload."),
    ("Conversation memory",    "Chat history enables follow-up questions."),
    ("Dynamic suggestions",    "UI suggests queries based on the dataset's actual columns."),
]
col2_x = LM + Inches(5.6)
for i, (title, desc) in enumerate(features):
    col = i // 3
    row = i % 3
    x = LM if col == 0 else col2_x
    item_y = y + row * Inches(1.2)
    accent_dot(s, x, item_y + Inches(0.06))
    txt(s, x + Inches(0.3), item_y, Inches(4.8), Inches(0.35),
        title, size=16, color=NEAR_BLACK, bold=True)
    txt(s, x + Inches(0.3), item_y + Inches(0.38), Inches(4.8), Inches(0.5),
        desc, size=13, color=GRAY)


# =====================================================================
# 8  EXAMPLES — text left, illustration right
# =====================================================================
s = prs.slides.add_slide(blank)
bg(s, OFF_WHITE)

txt(s, LM, TM, Inches(4), Inches(0.35),
    "Examples", size=13, color=BLUE, bold=True)

txt(s, LM, Inches(1.6), Inches(7), Inches(0.8),
    "Things you can ask.", size=30, color=NEAR_BLACK, bold=True)

examples = [
    ("Show the top 10 records by salary",               "Sorted query"),
    ("Find employees where department is Engineering",  "Filtered query"),
    ("Plot average age by city",                        "Bar chart"),
    ("Show count of records by status as pie chart",    "Pie chart"),
    ("What columns are in this dataset?",               "Schema info"),
]
y = Inches(2.9)
for query, kind in examples:
    thin_line(s, LM, y, Inches(7))
    txt(s, LM, y + Inches(0.2), Inches(5.8), Inches(0.4),
        f"\"{query}\"", size=15, color=NEAR_BLACK, italic=True)
    txt(s, Inches(7), y + Inches(0.22), Inches(1.5), Inches(0.35),
        kind, size=12, color=BLUE, align=PP_ALIGN.RIGHT)
    y += Inches(0.75)
thin_line(s, LM, y, Inches(7))

# Illustration on right
add_img(s, os.path.join(ASSETS, "demo.png"),
        Inches(8.8), Inches(1.2), width=Inches(4.0))


# =====================================================================
# 9  DEMO — dark break slide with illustration
# =====================================================================
s = prs.slides.add_slide(blank)
bg(s, DARK)

add_img(s, os.path.join(ASSETS, "demo.png"),
        Inches(1.0), Inches(1.0), height=Inches(5.5))

txt(s, Inches(6.5), Inches(2.6), Inches(6), Inches(0.9),
    "Demo", size=52, color=WHITE, bold=True)

thin_line(s, Inches(6.5), Inches(3.65), Inches(2.5), BLUE)

txt(s, Inches(6.5), Inches(4.0), Inches(6), Inches(0.5),
    "Let's see it in action.", size=20, color=LIGHT_GRAY)

# Demo steps
demo_y = Inches(4.8)
demo_steps = [
    "Upload a CSV (e.g., Heart Disease dataset)",
    "Ask: \"Show the top 5 records by age\"",
    "Ask: \"Plot average cholesterol by chest pain type\"",
    "Ask a follow-up question",
]
for step in demo_steps:
    txt(s, Inches(6.7), demo_y, Inches(5.5), Inches(0.35),
        f"→  {step}", size=13, color=LIGHT_GRAY)
    demo_y += Inches(0.4)


# =====================================================================
# 10  WHAT'S NEXT — text left, illustration right
# =====================================================================
s = prs.slides.add_slide(blank)
bg(s, WHITE)

txt(s, LM, TM, Inches(4), Inches(0.35),
    "What's next", size=13, color=BLUE, bold=True)

txt(s, LM, Inches(1.6), Inches(6.5), Inches(0.8),
    "Future improvements.", size=30, color=NEAR_BLACK, bold=True)

y = Inches(2.9)
future = [
    ("Multi-table support",  "Join across multiple uploaded files."),
    ("Authentication",       "Login, roles, and API key management."),
    ("Export results",       "Download charts as images, data as CSV."),
    ("RAG integration",      "Query documents alongside structured data."),
    ("Cloud deployment",     "Dockerized deploy on AWS/GCP."),
]
for i, (title, desc) in enumerate(future):
    numbered_item(s, LM, y, str(i + 1), title, desc, w=Inches(6))
    y += Inches(0.8)

# Illustration right
add_img(s, os.path.join(ASSETS, "future.png"),
        Inches(8.2), Inches(1.2), width=Inches(4.5))


# =====================================================================
# 11  THANK YOU
# =====================================================================
s = prs.slides.add_slide(blank)
bg(s, WHITE)

txt(s, Inches(0), Inches(2.4), SLIDE_W, Inches(0.9),
    "Thank you.", size=48, color=NEAR_BLACK, bold=True, align=PP_ALIGN.CENTER)

thin_line(s, Inches(5.5), Inches(3.5), Inches(2.3), BLUE)

txt(s, Inches(0), Inches(3.9), SLIDE_W, Inches(0.5),
    "Questions?", size=22, color=GRAY, align=PP_ALIGN.CENTER)

# Small hero image as a nice closer
add_img(s, os.path.join(ASSETS, "hero.png"),
        Inches(4.5), Inches(4.7), width=Inches(4.3))


# ── Save ─────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "Data_Chat_Agent_Presentation.pptx")
prs.save(out)
print(f"✅  Saved → {out}  ({len(prs.slides)} slides)")
